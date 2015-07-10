#!/usr/bin/env python

from __future__ import print_function
from docx import Document
from pptx import Presentation

from xml.sax.saxutils import escape, quoteattr
import os
from os.path import join, getsize

def printtag(name,closing=False):
    if closing:
        return '</{}>'.format(name)
    else:
        return '<{}>'.format(name)


def tag(name, *content):
    return '{0}{1}{2}'.format(printtag(name),
                             ''.join(content),
                             printtag(name, closing=True))

template = '''Path-Name: {0}
Content-Length: {1}
Document-Type: XML*

{2}'''

def docx2swishxml(fdocx):
    doc = Document(fdocx)
    content = ''.join([paragraph.text.encode('utf-8')
                       for paragraph in doc.paragraphs])
    content = ''.join([i if ord(i) < 128 else ' ' for i in content])
    xml = tag('swishdefault', escape(content))

    return template.format(fdocx, len(xml), xml)

def pptx2swishxml(pptx):
    prs = Presentation(pptx)
    content = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                content += shape.text.encode('utf-8')

    content = ''.join([i if ord(i) < 128 else ' ' for i in content])
    xml = tag('swishdefault', escape(content))

    return template.format(pptx, len(xml), xml)

index_dirs = ['/Users/jkitchin/Dropbox', '/Users/jkitchin/Box Sync']

for idir in index_dirs:
    for root, dirs, files in os.walk(idir):
        for f in files:
            try:
                # use end arg to avoid an extra space and carriage return
                fname = os.path.join(root, f)
                if fname.endswith('.docx'):
                    print(docx2swishxml(fname), end='')
                elif fname.endswith('.pptx'):
                    print(pptx2swishxml(fname), end='')
            except:
                pass
